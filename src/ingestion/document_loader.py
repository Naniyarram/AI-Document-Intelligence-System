
# src/ingestion/document_loader.py

# STAGE 1 — Smart Document Ingestion

# Handles: PDF, DOCX, DOC, TXT, XLSX, XLS, CSV, PNG, JPG
# For each page/section, detects what type of content it is:
#   - plain text
#   - table
#   - image / chart / diagram (sent to VLM later)
#   - scanned page (sent to OCR later)


import os
import io
from pathlib import Path
from dataclasses import dataclass, field
from typing import List, Optional
from loguru import logger


@dataclass
class DocumentPage:
    """
    One 'page' or 'section' from a document.
    This is the standard unit that flows through the pipeline.
    """
    # Where this content came from
    source_file: str          # e.g. "contract.pdf"
    page_number: int          # page index (0-based)
    section_title: str = ""   # heading above this content, if any

    # The actual content
    text: str = ""            # extracted text (empty if purely visual)

    # Content type — tells the pipeline how to handle this page
    # Options: "text", "table", "image", "scanned", "spreadsheet_row"
    content_type: str = "text"

    # For image/chart/diagram pages — raw image bytes for VLM
    image_bytes: Optional[bytes] = None
    image_format: str = "png"  # "png" or "jpeg"

    # Extra metadata (stored in vector DB alongside the text)
    metadata: dict = field(default_factory=dict)


class DocumentLoader:
    """
    Unified document loader.
    Detects file type and routes to the right parser.

    Usage:
        loader = DocumentLoader()
        pages = loader.load("my_contract.pdf")
        # returns List[DocumentPage]
    """

    def load(self, file_path: str) -> List[DocumentPage]:
        """
        Main entry point. Give it any supported file, get back a list of pages.
        """
        path = Path(file_path)
        extension = path.suffix.lower().lstrip(".")

        logger.info(f"Loading file: {path.name} (type: {extension})")

        # Route to the right parser based on file extension
        if extension == "pdf":
            return self._load_pdf(file_path)
        elif extension in ["docx", "doc"]:
            return self._load_docx(file_path)
        elif extension in ["xlsx", "xls"]:
            return self._load_excel(file_path)
        elif extension == "csv":
            return self._load_csv(file_path)
        elif extension == "txt":
            return self._load_txt(file_path)
        elif extension in ["png", "jpg", "jpeg"]:
            return self._load_image(file_path)
        elif extension == "pptx":
            return self._load_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: .{extension}")

   
    # PDF Parser
  
    def _load_pdf(self, file_path: str) -> List[DocumentPage]:
        """
        Parse PDF pages.

        For each page, we decide:
        - Does it have a text layer? → extract text
        - Does it have images/charts? → also capture image for VLM
        - Is it purely a scanned page? → mark as 'scanned' for OCR
        """
        import fitz  # PyMuPDF

        pages = []
        doc = fitz.open(file_path)
        filename = Path(file_path).name

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text from this page
            text = page.get_text("text").strip()

            # Check how many images are on this page
            image_list = page.get_images(full=True)
            has_images = len(image_list) > 0

            # Decide content type
            if not text and not has_images:
                # Empty page — skip it
                continue
            elif not text and has_images:
                # Purely a scanned page with no text layer
                content_type = "scanned"
                image_bytes = self._pdf_page_to_image(page)
                pages.append(DocumentPage(
                    source_file=filename,
                    page_number=page_num,
                    text="",
                    content_type=content_type,
                    image_bytes=image_bytes,
                    metadata={"page": page_num + 1, "source": filename}
                ))
            else:
                # Has text — check if it also has tables or embedded images
                # Check for table-like structure (lots of pipes or tab-separated text)
                if self._looks_like_table(text):
                    content_type = "table"
                else:
                    content_type = "text"

                # If the page ALSO has embedded charts/diagrams, capture image too
                image_bytes = None
                if has_images and len(text) < 200:
                    # Page is mostly image with little text — treat as image
                    content_type = "image"
                    image_bytes = self._pdf_page_to_image(page)
                elif has_images:
                    # Page has both text and images — capture image for VLM context
                    image_bytes = self._pdf_page_to_image(page)

                pages.append(DocumentPage(
                    source_file=filename,
                    page_number=page_num,
                    text=text,
                    content_type=content_type,
                    image_bytes=image_bytes,
                    metadata={"page": page_num + 1, "source": filename}
                ))

        doc.close()
        logger.info(f"PDF loaded: {len(pages)} pages from {filename}")
        return pages

    def _pdf_page_to_image(self, page) -> bytes:
        """Render a PDF page as a PNG image (for VLM processing)."""
        # Render at 150 DPI — good enough for VLM, not too heavy
        matrix = page._parent.get_page_pixmap(
            page.number,
            matrix=page._parent.get_page_pixmap.__module__
        ) if False else None

        # Simpler approach using fitz directly
        mat = __import__("fitz").Matrix(150 / 72, 150 / 72)  # 150 DPI
        pix = page.get_pixmap(matrix=mat)
        return pix.tobytes("png")

    # DOCX Parser
 
    def _load_docx(self, file_path: str) -> List[DocumentPage]:
        """
        Parse Word documents.
        Extracts paragraphs, tables, and detects section headings.
        """
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(file_path)
        filename = Path(file_path).name
        pages = []
        current_section = ""
        current_text_parts = []
        section_page_num = 0

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                # It's a paragraph
                para_text = "".join(
                    node.text for node in element.iter()
                    if node.tag.endswith("}t") and node.text
                ).strip()

                if not para_text:
                    continue

                # Detect if it's a heading (section title)
                style_name = ""
                pPr = element.find(qn("w:pPr"))
                if pPr is not None:
                    pStyle = pPr.find(qn("w:pStyle"))
                    if pStyle is not None:
                        style_name = pStyle.get(qn("w:val"), "")

                is_heading = "Heading" in style_name or "heading" in style_name.lower()

                if is_heading:
                    # Save accumulated text as a page before new section
                    if current_text_parts:
                        combined = "\n".join(current_text_parts)
                        pages.append(DocumentPage(
                            source_file=filename,
                            page_number=section_page_num,
                            section_title=current_section,
                            text=combined,
                            content_type="text",
                            metadata={"section": current_section, "source": filename}
                        ))
                        current_text_parts = []
                        section_page_num += 1

                    current_section = para_text
                else:
                    current_text_parts.append(para_text)

            elif tag == "tbl":
                # It's a table — extract as a text grid
                table_text = self._extract_docx_table(element)
                if table_text:
                    pages.append(DocumentPage(
                        source_file=filename,
                        page_number=section_page_num,
                        section_title=current_section,
                        text=table_text,
                        content_type="table",
                        metadata={"section": current_section, "source": filename, "is_table": True}
                    ))
                    section_page_num += 1

        # Don't forget the last accumulated text
        if current_text_parts:
            combined = "\n".join(current_text_parts)
            pages.append(DocumentPage(
                source_file=filename,
                page_number=section_page_num,
                section_title=current_section,
                text=combined,
                content_type="text",
                metadata={"section": current_section, "source": filename}
            ))

        logger.info(f"DOCX loaded: {len(pages)} sections from {filename}")
        return pages

    def _extract_docx_table(self, tbl_element) -> str:
        """Turn a DOCX table XML element into readable text rows."""
        from docx.oxml.ns import qn
        rows = []
        for row in tbl_element.findall(f".//{qn('w:tr')}"):
            cells = []
            for cell in row.findall(f".//{qn('w:tc')}"):
                cell_text = "".join(
                    node.text for node in cell.iter()
                    if node.tag.endswith("}t") and node.text
                ).strip()
                cells.append(cell_text)
            if any(cells):
                rows.append(" | ".join(cells))
        return "\n".join(rows)

    # Excel Parser

    def _load_excel(self, file_path: str) -> List[DocumentPage]:
        """
        Parse Excel files.
        Each sheet becomes a collection of pages.
        Large sheets are chunked into groups of rows.
        """
        import openpyxl

        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        filename = Path(file_path).name
        pages = []
        page_num = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text = []
            headers = []

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                # Skip completely empty rows
                row_values = [str(v) if v is not None else "" for v in row]
                if not any(v.strip() for v in row_values):
                    continue

                if row_idx == 0:
                    # First row is usually headers
                    headers = row_values
                    rows_text.append("COLUMNS: " + " | ".join(headers))
                else:
                    # Format as "Header: Value" pairs for better readability
                    if headers:
                        pairs = [f"{h}: {v}" for h, v in zip(headers, row_values) if h or v]
                        rows_text.append(" | ".join(pairs))
                    else:
                        rows_text.append(" | ".join(row_values))

                # Group every 30 rows into one page (prevents huge chunks)
                if len(rows_text) >= 30:
                    page_text = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
                    pages.append(DocumentPage(
                        source_file=filename,
                        page_number=page_num,
                        section_title=f"Sheet: {sheet_name}",
                        text=page_text,
                        content_type="spreadsheet_row",
                        metadata={"sheet": sheet_name, "source": filename, "is_table": True}
                    ))
                    page_num += 1
                    rows_text = []
                    if headers:
                        rows_text.append("COLUMNS: " + " | ".join(headers))  # re-add headers

            # Remaining rows
            if rows_text:
                page_text = f"Sheet: {sheet_name}\n" + "\n".join(rows_text)
                pages.append(DocumentPage(
                    source_file=filename,
                    page_number=page_num,
                    section_title=f"Sheet: {sheet_name}",
                    text=page_text,
                    content_type="spreadsheet_row",
                    metadata={"sheet": sheet_name, "source": filename, "is_table": True}
                ))
                page_num += 1

        wb.close()
        logger.info(f"Excel loaded: {len(pages)} sections from {filename}")
        return pages


    # CSV Parser

    def _load_csv(self, file_path: str) -> List[DocumentPage]:
        """Parse CSV files as tabular data."""
        import csv

        filename = Path(file_path).name
        pages = []
        rows_text = []
        headers = []
        page_num = 0

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row_idx, row in enumerate(reader):
                if row_idx == 0:
                    headers = row
                    rows_text.append("COLUMNS: " + " | ".join(headers))
                else:
                    if headers:
                        pairs = [f"{h}: {v}" for h, v in zip(headers, row) if h or v]
                        rows_text.append(" | ".join(pairs))
                    else:
                        rows_text.append(" | ".join(row))

                if len(rows_text) >= 40:
                    pages.append(DocumentPage(
                        source_file=filename,
                        page_number=page_num,
                        text="\n".join(rows_text),
                        content_type="spreadsheet_row",
                        metadata={"source": filename, "is_table": True}
                    ))
                    page_num += 1
                    rows_text = []
                    if headers:
                        rows_text.append("COLUMNS: " + " | ".join(headers))

        if rows_text:
            pages.append(DocumentPage(
                source_file=filename,
                page_number=page_num,
                text="\n".join(rows_text),
                content_type="spreadsheet_row",
                metadata={"source": filename, "is_table": True}
            ))

        return pages

  
    # TXT Parser

    def _load_txt(self, file_path: str) -> List[DocumentPage]:
        """Parse plain text files. Split by double newlines into sections."""
        filename = Path(file_path).name

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()

        # Split on double newlines (natural paragraph breaks)
        sections = [s.strip() for s in content.split("\n\n") if s.strip()]

        pages = []
        for idx, section in enumerate(sections):
            pages.append(DocumentPage(
                source_file=filename,
                page_number=idx,
                text=section,
                content_type="text",
                metadata={"source": filename, "section_idx": idx}
            ))

        return pages


    # Image Parser

    def _load_image(self, file_path: str) -> List[DocumentPage]:
        """
        Load an image file (PNG/JPG).
        Marks it for VLM processing — no text extraction here.
        The VLM module will handle what's in the image.
        """
        filename = Path(file_path).name

        with open(file_path, "rb") as f:
            image_bytes = f.read()

        extension = Path(file_path).suffix.lower().lstrip(".")
        img_format = "jpeg" if extension in ["jpg", "jpeg"] else "png"

        return [DocumentPage(
            source_file=filename,
            page_number=0,
            text="",
            content_type="image",
            image_bytes=image_bytes,
            image_format=img_format,
            metadata={"source": filename}
        )]


    # PPTX Parser
 
    def _load_pptx(self, file_path: str) -> List[DocumentPage]:
        """Parse PowerPoint files slide by slide."""
        from pptx import Presentation

        prs = Presentation(file_path)
        filename = Path(file_path).name
        pages = []

        for slide_num, slide in enumerate(prs.slides):
            text_parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            text_parts.append(text)

            if text_parts:
                pages.append(DocumentPage(
                    source_file=filename,
                    page_number=slide_num,
                    section_title=f"Slide {slide_num + 1}",
                    text="\n".join(text_parts),
                    content_type="text",
                    metadata={"slide": slide_num + 1, "source": filename}
                ))

        return pages

    # Helper Methods
 
    def _looks_like_table(self, text: str) -> bool:
        """
        Quick heuristic: does this text look like a table?
        Checks for pipe characters or many tab-separated columns.
        """
        lines = text.split("\n")
        pipe_lines = sum(1 for line in lines if "|" in line)
        tab_lines = sum(1 for line in lines if "\t" in line and len(line.split("\t")) > 2)
        total_lines = len(lines)

        if total_lines == 0:
            return False

        # If more than 30% of lines look table-like, call it a table
        return (pipe_lines / total_lines > 0.3) or (tab_lines / total_lines > 0.3)
